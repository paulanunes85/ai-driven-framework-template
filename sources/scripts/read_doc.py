#!/usr/bin/env python3
"""Unified document reader for the framework workspace.

Extracts structured text from any supported file type.
Supports: PPTX, DOCX, XLSX, PDF, MD, VTT

Usage:
    python sources/scripts/read_doc.py <file>
    python sources/scripts/read_doc.py <file> --json
"""

import sys
import os
import json
import re
from pathlib import Path


# ---------------------------------------------------------------------------
# Dependency checks (conditional imports with graceful fallback)
# ---------------------------------------------------------------------------
_MISSING_DEPS: dict[str, str] = {}

try:
    from pptx import Presentation as _Presentation  # noqa: F401
except ImportError:
    _MISSING_DEPS[".pptx"] = "python-pptx"

try:
    from docx import Document as _Document  # noqa: F401
except ImportError:
    _MISSING_DEPS[".docx"] = "python-docx"

try:
    from openpyxl import load_workbook as _load_workbook  # noqa: F401
except ImportError:
    _MISSING_DEPS[".xlsx"] = "openpyxl"

try:
    import fitz as _fitz  # noqa: F401
except ImportError:
    _MISSING_DEPS[".pdf"] = "pymupdf"


def _check_dep(ext: str) -> None:
    """Raise a helpful error if a required dependency is missing."""
    if ext in _MISSING_DEPS:
        pkg = _MISSING_DEPS[ext]
        raise ImportError(
            f"Reading {ext} files requires '{pkg}'. "
            f"Install it with:  pip install {pkg}"
        )


# ---------------------------------------------------------------------------
# PPTX Reader
# ---------------------------------------------------------------------------
def read_pptx(path: str) -> dict:
    _check_dep(".pptx")
    from pptx import Presentation

    prs = Presentation(path)
    result = {
        "type": "pptx",
        "file": os.path.basename(path),
        "slide_width": str(prs.slide_width),
        "slide_height": str(prs.slide_height),
        "slide_count": len(prs.slides),
        "slides": [],
    }

    for i, slide in enumerate(prs.slides, 1):
        slide_data = {
            "number": i,
            "layout": slide.slide_layout.name if slide.slide_layout else None,
            "shapes": [],
            "notes": None,
        }

        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "type": (
                    shape.shape_type.__class__.__name__
                    if hasattr(shape.shape_type, "__class__")
                    else str(shape.shape_type)
                ),
                "position": {
                    "left": str(shape.left),
                    "top": str(shape.top),
                    "width": str(shape.width),
                    "height": str(shape.height),
                },
            }

            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        para_info = {"text": text}
                        if para.runs:
                            run = para.runs[0]
                            if run.font.size:
                                para_info["font_size"] = str(run.font.size)
                            if run.font.bold:
                                para_info["bold"] = True
                        paragraphs.append(para_info)
                if paragraphs:
                    shape_info["paragraphs"] = paragraphs

            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                shape_info["table"] = table_data

            slide_data["shapes"].append(shape_info)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data["notes"] = notes_text

        result["slides"].append(slide_data)

    return result


# ---------------------------------------------------------------------------
# DOCX Reader
# ---------------------------------------------------------------------------
def read_docx(path: str) -> dict:
    _check_dep(".docx")
    from docx import Document

    doc = Document(path)
    result = {
        "type": "docx",
        "file": os.path.basename(path),
        "properties": {},
        "sections": [],
        "tables": [],
    }

    props = doc.core_properties
    for attr in [
        "title", "subject", "author", "category",
        "comments", "keywords", "created", "modified",
        "last_modified_by", "revision",
    ]:
        val = getattr(props, attr, None)
        if val:
            result["properties"][attr] = str(val)

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        entry = {
            "text": text,
            "style": para.style.name if para.style else None,
        }
        result["sections"].append(entry)

    for i, table in enumerate(doc.tables):
        table_data = {"index": i, "rows": []}
        for row in table.rows:
            table_data["rows"].append(
                [cell.text.strip() for cell in row.cells]
            )
        result["tables"].append(table_data)

    return result


# ---------------------------------------------------------------------------
# XLSX Reader
# ---------------------------------------------------------------------------
def read_xlsx(path: str) -> dict:
    _check_dep(".xlsx")
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True)
    result = {
        "type": "xlsx",
        "file": os.path.basename(path),
        "properties": {},
        "sheets": [],
    }

    props = wb.properties
    for attr in ["title", "subject", "creator", "category", "description",
                 "created", "modified", "lastModifiedBy"]:
        val = getattr(props, attr, None)
        if val:
            result["properties"][attr] = str(val)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_data = {
            "name": sheet_name,
            "dimensions": ws.dimensions,
            "rows": [],
        }
        for row in ws.iter_rows(values_only=True):
            row_vals = []
            for cell in row:
                if cell is not None:
                    row_vals.append(str(cell))
                else:
                    row_vals.append("")
            if any(row_vals):
                sheet_data["rows"].append(row_vals)
        result["sheets"].append(sheet_data)

    wb.close()
    return result


# ---------------------------------------------------------------------------
# Markdown Reader
# ---------------------------------------------------------------------------
def read_md(path: str) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        content = f.read()

    headings = []
    for match in re.finditer(r'^(#{1,6})\s+(.+)$', content, re.MULTILINE):
        headings.append({
            "level": len(match.group(1)),
            "text": match.group(2).strip(),
        })

    tables = []
    table_pattern = re.compile(
        r'(\|.+\|)\n(\|[-| :]+\|)\n((?:\|.+\|\n?)+)',
        re.MULTILINE,
    )
    for match in table_pattern.finditer(content):
        header = [c.strip() for c in match.group(1).split("|") if c.strip()]
        rows = []
        for line in match.group(3).strip().split("\n"):
            row = [c.strip() for c in line.split("|") if c.strip()]
            rows.append(row)
        tables.append({"header": header, "rows": rows})

    links = []
    for match in re.finditer(r'\[([^\]]+)\]\(([^)]+)\)', content):
        links.append({"text": match.group(1), "url": match.group(2)})

    return {
        "type": "md",
        "file": os.path.basename(path),
        "line_count": content.count("\n") + 1,
        "char_count": len(content),
        "headings": headings,
        "tables": tables,
        "internal_links": [
            lnk for lnk in links
            if not lnk["url"].startswith("http")
        ],
    }


# ---------------------------------------------------------------------------
# PDF Reader
# ---------------------------------------------------------------------------
def read_pdf(path: str) -> dict:
    _check_dep(".pdf")
    import fitz  # PyMuPDF

    doc = fitz.open(path)
    result = {
        "type": "pdf",
        "file": os.path.basename(path),
        "page_count": len(doc),
        "metadata": {},
        "pages": [],
    }

    meta = doc.metadata or {}
    for key in [
        "title", "author", "subject", "creator",
        "producer", "creationDate", "modDate",
    ]:
        val = meta.get(key, "")
        if val:
            result["metadata"][key] = val

    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        result["pages"].append({
            "number": i,
            "width": round(page.rect.width, 1),
            "height": round(page.rect.height, 1),
            "char_count": len(text),
            "text": text,
        })

    doc.close()
    return result


# ---------------------------------------------------------------------------
# VTT Reader
# ---------------------------------------------------------------------------
def read_vtt(path: str) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        content = f.read()

    lines = content.split("\n")
    cues = []
    speakers: set[str] = set()
    current_time = None

    for line in lines:
        line = line.strip()
        # Timestamp line: 00:00:00.000 --> 00:00:05.000
        if "-->" in line:
            current_time = line.split("-->")[0].strip()
            continue
        if (
            line
            and not line.startswith("WEBVTT")
            and not line.startswith("NOTE")
            and current_time
        ):
            # Extract speaker if present: "Speaker Name: text"
            speaker_match = re.match(r'^([^:]{1,50}):\s*(.+)', line)
            if speaker_match:
                speaker = speaker_match.group(1).strip()
                speakers.add(speaker)
            cues.append({"time": current_time, "text": line})
            current_time = None

    return {
        "type": "vtt",
        "file": os.path.basename(path),
        "cue_count": len(cues),
        "speakers": sorted(speakers),
        "duration_approx": cues[-1]["time"] if cues else "00:00:00.000",
        "cues_sample": cues[:10],
    }


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------
READERS = {
    ".pptx": read_pptx,
    ".docx": read_docx,
    ".xlsx": read_xlsx,
    ".md": read_md,
    ".pdf": read_pdf,
    ".vtt": read_vtt,
}


def read_file(fpath: str) -> dict:
    """Read a file and return structured data."""
    ext = Path(fpath).suffix.lower()
    reader = READERS.get(ext)
    if not reader:
        return {
            "error": f"Unsupported file type: {ext}",
            "supported": list(READERS.keys()),
        }
    return reader(fpath)


def format_text(result: dict) -> str:
    """Pretty-print extracted data as human-readable text."""
    lines = []
    ftype = result.get("type", "unknown")
    lines.append(f"{'=' * 70}")
    lines.append(f"  {result.get('file', '?')}  ({ftype.upper()})")
    lines.append(f"{'=' * 70}")

    if ftype == "pptx":
        lines.append(f"Slides: {result['slide_count']}")
        for slide in result["slides"]:
            lines.append(f"\n--- Slide {slide['number']} ---")
            for shape in slide["shapes"]:
                if "paragraphs" in shape:
                    for p in shape["paragraphs"]:
                        prefix = "  [B] " if p.get("bold") else "      "
                        lines.append(f"{prefix}{p['text']}")
                if "table" in shape:
                    for row in shape["table"]:
                        lines.append(f"  | {'  |  '.join(row)} |")
            if slide.get("notes"):
                lines.append(f"  Notes: {slide['notes'][:100]}...")

    elif ftype == "docx":
        if result.get("properties"):
            for k, v in result["properties"].items():
                lines.append(f"  {k}: {v}")
        lines.append("")
        for sec in result["sections"]:
            style = sec.get("style", "")
            if "Heading" in style:
                level = (
                    int(style[-1]) if style[-1].isdigit()
                    else 1
                )
                lines.append(f"\n{'#' * level} {sec['text']}")
            else:
                lines.append(f"  {sec['text']}")

    elif ftype == "xlsx":
        if result.get("properties"):
            for k, v in result["properties"].items():
                lines.append(f"  {k}: {v}")
        for sheet in result["sheets"]:
            dims = sheet["dimensions"]
            lines.append(
                f"\n--- Sheet: {sheet['name']}"
                f" ({dims}) ---"
            )
            for row in sheet["rows"]:
                lines.append(f"  | {'  |  '.join(row)} |")

    elif ftype == "md":
        lines.append(
            f"Lines: {result['line_count']}"
            f"  |  Chars: {result['char_count']}"
        )
        lines.append(f"Headings: {len(result['headings'])}")
        for h in result["headings"]:
            lines.append(f"  {'#' * h['level']} {h['text']}")
        if result["tables"]:
            lines.append(f"Tables: {len(result['tables'])}")
        if result["internal_links"]:
            lines.append(
                f"Internal links: {len(result['internal_links'])}"
            )

    elif ftype == "pdf":
        if result.get("metadata"):
            for k, v in result["metadata"].items():
                lines.append(f"  {k}: {v}")
        lines.append(f"Pages: {result['page_count']}")
        for page in result["pages"]:
            lines.append(
                f"\n--- Page {page['number']}"
                f" ({page['char_count']} chars) ---"
            )
            text = page["text"]
            if len(text) > 500:
                lines.append(f"  {text[:500]}...")
            else:
                lines.append(f"  {text}")

    elif ftype == "vtt":
        lines.append(f"Cues: {result['cue_count']}")
        lines.append(f"Duration: ~{result['duration_approx']}")
        if result.get("speakers"):
            lines.append(
                f"Speakers: {', '.join(result['speakers'])}"
            )
        if result.get("cues_sample"):
            lines.append(
                f"\nFirst {len(result['cues_sample'])} cues:"
            )
            for cue in result["cues_sample"]:
                lines.append(f"  [{cue['time']}] {cue['text'][:80]}")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__.strip())
        sys.exit(1)

    filepath = sys.argv[1]
    if not os.path.isfile(filepath):
        print(f"Error: file not found -- {filepath}")
        sys.exit(1)

    as_json = "--json" in sys.argv

    try:
        data = read_file(filepath)
    except ImportError as exc:
        print(f"Error: {exc}")
        sys.exit(1)

    if "error" in data:
        print(f"Error: {data['error']}")
        print(f"Supported types: {', '.join(data['supported'])}")
        sys.exit(1)

    if as_json:
        print(json.dumps(data, indent=2, ensure_ascii=False, default=str))
    else:
        print(format_text(data))
